"""
This file creates a PowerPoint presentation with plots and metrics from the results of multiple runs.
It does not create the plots itself, but assumes that they are already available in the result directory of each individual run.
The settings that you need to change are marked as TODO.
"""

import dotenv
import glob
import os
from pathlib import Path
import pptx


def slide_run(d, prs):
    slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(slide_layout)
    left = pptx.util.Cm(0.5)
    top = pptx.util.Cm(2)
    high = pptx.util.Cm(5)  # height of standard plot
    add = pptx.util.Cm(0.5)  # vertical space between plots
    col_width = pptx.util.Cm(11)  # column width

    # Set directory path as the slide title
    title1 = slide.shapes.title
    title1.width = pptx.util.Cm(24)
    title1.height = pptx.util.Cm(1)
    title1.text = str(d)
    title1.text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.LEFT
    title1.text_frame.paragraphs[0].font.size = pptx.util.Pt(12)  # adjust the font size

    # Plots; name, vertical position, and height of the plots; each list is one column
    plots_and_settings = [[("duration", top, high / 3 * 2),
                           ("mean_daily", top + high + add, high)],
                          [("p_infl", top, high),
                           ("price_n", top + high + add, high),
                           ("p_tot", top + 2 * high + 2 * add, high)],
                          [("p_wh", top, high),
                           ("p_hp", top + high + add, high),
                           ("p_ev", top + 2 * high + 2 * add, high)],
                          [("highest_p_tot", top, high),
                           ("lowest_p_infl", top + high + add, high),
                           ("highest_p_infl", top + 2 * high + 2 * add, high)]
                          ]
    for i, col in enumerate(plots_and_settings):
        for name, pos_y, height in col:
            # Add plot name above each figure
            pos_x = left + i * col_width
            textbox = slide.shapes.add_textbox(pos_x, pos_y - pptx.util.Cm(0.5), width=col_width - pptx.util.Cm(1), height=pptx.util.Cm(0.5))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = pptx.util.Pt(9)
            p.alignment = pptx.enum.text.PP_ALIGN.LEFT

            # Add figure below the plot name
            file = glob.glob(os.path.join(d, f"*{name}*.png"))[0]
            slide.shapes.add_picture(file, pos_x, pos_y, height=height)

    # Set slide width
    prs.slide_width = len(plots_and_settings) * col_width + 2 * left

    # Add metrics from metrics.txt file and first part of log (config) as notes
    with open(d / "metrics.txt", "r") as file:
       metrics_text = file.read()
    with open(d.parent / "log_file.log", "r") as file:
        lines = file.readlines()
        config = ""
        for line in lines:
            if line.strip() == "":
                break
            config += line
    slide.notes_slide.notes_text_frame.text = metrics_text + "\n\n" + config

    return prs


def main():
    # Settings
    path_res = Path(os.environ["PATH_SIM_RESULTS"])  # results directory
    list_dirs = [  # list of runs to be included in the presentation
        "240302_172421_rl_dynamic/final_eval",  # TODO: replace with actual directories
        "240302_172434_rl_dynamic/final_eval",
    ]
    path_output = Path(os.environ["PATH_ANALYSIS_RESULTS"])  # output directory
    file_output = "sweep_59bj8cpp.pptx"  # file name of the presentation TODO: replace with desired file name

    # Initialize presentation
    prs = pptx.Presentation()

    # Individual slide for each run
    for d in list_dirs:
        prs = slide_run(path_res / d, prs)

    prs.save(path_output / file_output)


if __name__ == "__main__":
    # Environment variables
    dotenv.load_dotenv(Path("../config/local.env"))

    main()
